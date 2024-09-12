import streamlit as st
import google.generativeai as genai
import os
import pdfplumber
import base64
from io import BytesIO
from PIL import Image
from dotenv import load_dotenv, set_key
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
import openpyxl
import markdown

# 加载.env文件
load_dotenv()

# 初始化会话状态变量
for key in ['analysis_result', 'combined_content']:
    if key not in st.session_state:
        st.session_state[key] = None

def get_api_key():
    api_key = st.secrets.get("GOOGLE_API_KEY")
    if not api_key:
        api_key = st.sidebar.text_input("输入您的Google Gemini API密钥", type="password")
        if api_key:
            st.sidebar.success("API密钥已设置并保存!")
    else:
        st.sidebar.success("已从Secrets中读取API密钥")
    return api_key

st.title("文档分析器")

api_key = get_api_key()
max_tokens = st.slider("选择要处理的最大文本长度 (token)", 1000, 1000000, 128000)
temperature = st.slider("选择模型的temperature", 0.0, 1.0, 0.3, 0.1)

def process_image(image, caption):
    buffered = BytesIO()
    image.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return {'type': 'image', 'data': img_str, 'caption': caption}

def extract_content_from_pdf(pdf_file):
    with pdfplumber.open(pdf_file) as pdf:
        content = []
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                content.append(f"第 {page_num} 页文本:\n{text}\n")
            
            for i, table in enumerate(page.extract_tables(), 1):
                content.append(f"第 {page_num} 页表格 {i}:")
                content.extend(" | ".join(str(cell) for cell in row) for row in table)
                content.append("")
            
            for i, img in enumerate(page.images, 1):
                try:
                    image = Image.open(BytesIO(img['stream'].get_data()))
                    content.append(process_image(image, f"第 {page_num} 页图像 {i}"))
                except Exception as e:
                    content.append(f"第 {page_num} 页图像 {i}: [处理图像时出错: {str(e)}]")
        return content

def extract_content_from_pptx(file):
    content = []
    prs = Presentation(file)
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape_num, shape in enumerate(slide.shapes, 1):
            if shape.has_text_frame:
                slide_content.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_content = [f"幻灯片 {slide_num} 表格 {shape_num}:"]
                table_content.extend(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows)
                slide_content.extend(table_content)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = Image.open(BytesIO(shape.image.blob))
                content.append(process_image(image, f"幻灯片 {slide_num} 图像 {shape_num}"))
        content.append(f"幻灯片 {slide_num} 内容:\n" + "\n".join(slide_content))
    return content

def iter_block_items(parent):
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("不支持的父元素类型")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)

def extract_content_from_docx(file):
    doc = Document(file)
    content = []
    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            content.append(block.text)
        elif isinstance(block, Table):
            table_content = []
            for row in block.rows:
                table_content.append(" | ".join(cell.text for cell in row.cells))
            content.append("\n".join(table_content))

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image = Image.open(BytesIO(rel.target_part.blob))
            content.append(process_image(image, "文档图像"))
    return content

def extract_content_from_xlsx(file):
    wb = openpyxl.load_workbook(file)
    content = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        content.append(f"工作表: {sheet}")
        content.extend(" | ".join(str(cell.value) for cell in row) for row in ws.iter_rows())
        for image in ws._images:
            img = Image.open(BytesIO(image.ref))
            content.append(process_image(img, f"工作表 {sheet} 图像"))
    return content

def extract_content_from_text(file):
    content = file.getvalue().decode('utf-8')
    return {'type': 'text', 'data': markdown.markdown(content) if file.name.endswith('.md') else content}

def extract_content_from_image(file):
    try:
        return process_image(Image.open(file), "上传的图像")
    except Exception as e:
        return {'type': 'error', 'data': f"处理图像时出错: {str(e)}"}

def extract_content_from_file(file):
    extractors = {
        '.pdf': extract_content_from_pdf,
        '.pptx': extract_content_from_pptx,
        '.docx': extract_content_from_docx,
        '.xlsx': extract_content_from_xlsx,
        '.txt': extract_content_from_text,
        '.md': extract_content_from_text,
        '.jpg': extract_content_from_image,
        '.jpeg': extract_content_from_image,
        '.png': extract_content_from_image,
        '.gif': extract_content_from_image,
        '.bmp': extract_content_from_image
    }
    file_extension = os.path.splitext(file.name)[1].lower()
    extractor = extractors.get(file_extension)
    return extractor(file) if extractor else {'type': 'error', 'data': f"不支持的文件类型: {file_extension}"}

def get_prompt_options():
    return {
        "选项1 - 综合分析": """
        请分析以下多个文件的内容, 包括文本、表格和图像描述，并按步骤进行推理，逐步得出结论，提供系统化的总结。请确保在每个步骤中清晰地列出推理过程中的中间结论和逻辑:

        {combined_content}

        1. **所有文档的主要主题**: 请首先梳理各个文件的主要主题，按顺序列出每个文档的主题，并提供推理依据，说明为什么这些主题是最为重要的。

        2. **关键要点 (列出5-7个要点，包括文本、表格和图像中的信息)**: 请逐个文件分析关键要点，先对每个文件内容进行初步分析，然后根据文本、表格、图像内容得出5-7个关键要点。列出每个要点时，请说明你是如何从内容推导出这些要点的。
        
        3. **主要结论**: 在列出关键要点之后，进行总结，推导出每个文件的主要结论。请简要说明是如何通过前面的分析过程得出这些结论的，并为每个结论提供推理依据。
        
        4. **建议的后续行动 (如果适用)**: 基于主要结论，分析是否有必要采取进一步行动。如果有，列出建议的后续行动，并说明这些建议的推理依据或原因。
        
        5. **文档之间的关联性或差异 (如果适用)**: 分析多个文件之间的关联性或差异。在得出最终结论之前，请逐步进行比较，列出关联点或不同点，并简要说明你是如何得出这些相似或差异之处的。
        
        6. **表格和图像的重要信息摘要**: 对表格和图像中的信息进行专门的分析，按顺序列出图表中的重要数据或图像描述的重要内容。在总结之前，请说明你是如何从这些表格和图像中提取出关键信息的，并简述它们的意义。
        
        通过每一步的推理，确保总结的内容具有连贯性和逻辑性。
        """,
        "选项2 - 创建速查表": """
        基于以下多个文件的内容，创建一个简洁的速查表，并通过逐步推理确保速查表中的每一项信息具有逻辑性和清晰性。请在每一步中解释信息的来源及其重要性，以确保速查表简明且实用:

        {combined_content}

        速查表应包含以下内容，并按以下步骤推理完成:
        
        1. **主要概念和定义**: 逐步分析文件中的核心概念和定义。请先梳理每个文件中的关键术语，然后解释为什么这些术语是速查表中必须包含的。列出定义时，请简要说明这些概念如何帮助用户快速理解文档内容。
        
        2. **关键流程或步骤**: 从多个文件中提取关键流程或步骤。请先分析每个文件的流程，并解释每一步流程的意义。通过推理确定哪些步骤最为关键，并确保流程顺序清晰且易于参考。列出时请说明你如何从文件中推导出这些关键步骤。
        
        3. **重要数据或统计信息**: 请首先在每个文件中确定数据或统计信息，然后逐步解释这些信息的重要性。列出重要数据时，请说明数据的来源及其背后的逻辑，并解释这些数据对理解整体内容的帮助。
        
        4. **最佳实践或建议**: 结合文件内容，推导出有助于实践的建议或最佳做法。请先分析文件中的建议部分，逐步推理得出最佳实践，确保这些建议是基于文件内容而得出的，并具有实际应用价值。
        
        5. **常见问题及解决方案**: 从文件中推导出可能的常见问题及其解决方案。请逐步分析文件中的问题和对应的解决方案，并通过推理过程总结出对用户最有帮助的内容，确保这些问题和解决方案具有实用性和针对性。
        
        请确保所有信息都以易于理解和快速参考的格式呈现，并在每个部分的推理过程中简要解释如何从文档中得出这些内容。
        """,
        "选项3 - PDCA分析": """
        使用PDCA (计划-执行-检查-行动) 方法分析以下多个文件中的问题，并通过逐步推理制定行动计划。请确保在每一步中清晰说明问题的识别、分析过程及行动建议的逻辑依据:

        {combined_content}

        请按照以下步骤提供内容:
        1. **问题识别**: 逐步分析每个文档中的内容，识别主要问题或挑战。请为每个问题提供简要说明，并解释你是如何从文档中推导出这些问题的。
        
        2. **严重性分析**: 对每个问题进行严重性评估，逐步推理出问题的影响程度（高/中/低）。请解释为什么某些问题被评为高严重性，而其他问题被评为中或低，并说明每个问题可能带来的后果。
        
        3. **根本原因分析**: 对识别出的每个问题，进行根本原因分析。请详细说明如何通过文档中的信息推导出潜在的根本原因，并解释这些原因是如何导致问题发生的。
        
        4. **行动计划**: 
            - **针对高严重性问题**: 提供即时行动建议，并说明这些建议如何帮助迅速解决或缓解问题。请确保建议有逻辑性，且与前面的分析一致。
            - **针对中等严重性问题**: 提供中期行动建议，逐步推理出需要采取的措施，解释这些措施的实施时机和预期效果。
            - **针对低严重性问题**: 提供长期改进建议，详细说明这些问题的长期影响，以及如何逐步改善或解决。

        5. **监控和评估建议**: 提出如何跟踪和衡量改进措施的效果。请逐步推理出最佳的监控方法，并解释为什么这些方法能有效评估每个问题的改进效果。

        通过每一步的分析，确保所有问题的识别、评估和行动计划具有清晰的推理逻辑和可执行性。
        """,
        "选项4 - PowerPoint演示计划": """
        基于以下文档内容，创建一个详细且经过逐步推理的PowerPoint演示文稿计划。请确保每个步骤中的内容选择有逻辑依据，并逐步推导出演示文稿的结构和要点:

        {combined_content}

        请按照以下格式提供PowerPoint演示计划，并在每一步中解释内容选择的理由:

        1. **标题幻灯片**:
            - **主标题**: [请根据文档的核心主题推导出适当的主标题，并解释为什么选择这个标题]
            - **副标题**: [请简要说明副标题的作用，确保它能补充或进一步阐明主标题]

        2. **目录幻灯片**:
            - [列出演示文稿的主要部分，并解释每个部分在演示中占据的重要性。请逐步分析文档内容，推导出需要涵盖的关键部分]

        3. **内容幻灯片 (为每个主要部分创建2-3张幻灯片)**:
            a. **[第一部分标题]**:
                - **要点1**: [简洁描述主要要点，并解释为什么这个要点对观众理解该部分内容至关重要]
                - **要点2**: [继续描述，并说明要点的推导依据]
                - **要点3**: [说明要点的选择和它对整个内容的贡献]
                [**建议的视觉元素**: 请解释为什么图表、图像或图标能有效传达该部分的关键信息]

            b. **[第二部分标题]**:
                - **要点1**: [简洁描述，并说明如何从文档中得出这一要点]
                - **要点2**: [简述要点及其在演示中的作用]
                - **要点3**: [解释其选择的理由及其对整体内容的贡献]
                [**建议的视觉元素**: 说明选择这些视觉元素的原因以及它们如何增强演示效果]

            c. **[继续添加更多部分...]**

        4. **关键数据幻灯片**:
            - [列出需要突出的重要数据或统计信息，并解释这些数据如何支持演示内容的关键论点]
            [**建议的视觉元素**: 说明选择图表、表格或信息图的理由，确保它们能清晰有效地展示数据]

        5. **总结幻灯片**:
            - **主要结论1**: [简述主要结论，并解释如何从之前的内容推导出该结论]
            - **主要结论2**: [简述并说明该结论的重要性]
            - **主要结论3**: [确保每个结论与演示内容逻辑一致，并简述推理过程]

        6. **行动步骤幻灯片**:
            - **建议的后续行动1**: [简述后续行动的必要性，并解释如何从文档内容中推导出这些建议]
            - **建议的后续行动2**: [简述推导依据，并说明其执行的可行性]
            - **建议的后续行动3**: [提供逻辑推理过程，以确保行动步骤具有清晰的实施路径]

        7. **结束幻灯片**:
            - **感谢语**: [请根据演示的内容简要总结结束语，并确保它与观众的期望相符]
            - **联系信息或下一步指示**: [确保这些信息清晰、简洁，并对后续行动有指导性作用]

        注意:
            - 每张幻灯片的内容应简洁、重点突出，并有逻辑推导的支持
            - 使用简明的语言和短句，以确保易于理解
            - 视觉元素需与内容匹配，增强观众对信息的吸收
            - 整个演示文稿应保持一致的主题和风格，并通过逐步推理使内容结构清晰
        """,
        "选项5 - 图像分析": """
        分析以下文件内容，包括文本、表格、图像和其他数据。请按步骤推理，通过逐步分析从各类内容中得出清晰的结论:

        {combined_content}

        对于**图像内容**，请按以下步骤逐步分析:
            1. **描述你看到的内容**: 请详细描述图像中的元素，解释你是如何确定每个元素的，并说明其在图像中的位置和作用。
            2. **识别主要对象、场景或主题**: 通过分析图像中的视觉线索，识别图像中的主要对象、场景或主题，并简述推理过程，解释你为什么认为这些元素是图像的焦点。
            3. **分析图像的整体氛围或情感**: 根据图像中的色彩、构图和主题，分析其传达的氛围或情感。请解释你是如何从图像的视觉元素中推导出这些情感或氛围的。
            4. **提出与图像相关的任何见解或观察**: 基于前面的分析，提供你的个人见解或观察。请说明你是如何得出这些见解的，以及它们与图像中的具体元素之间的关系。

        接下来，请提供以下格式的总结，并通过逐步推理阐述每个结论:
            1. **文本内容的主要主题和关键点**: 从文本部分分析出主要主题和关键信息，并解释如何从文档中得出这些要点。
            2. **图像内容的描述和分析**: 总结图像的描述，并简要重述你对图像的分析结果。请确保你在总结时，解释每个视觉分析步骤的推理过程。
            3. **文本和图像之间的任何关联或互补信息**: 通过比较文本和图像的内容，分析它们之间的联系或互补性。请说明你是如何发现这些关联的，并解释这些联系如何加强了文件的整体理解。
            4. **整体结论和见解**: 基于所有分析，给出你的整体结论。请解释你是如何从不同类型的内容中推导出这个结论的，并总结出主要的洞察。
            5. **基于所有内容的建议或后续行动**: 根据整体分析结果，提供可行的建议或后续行动，并解释你为什么会提出这些建议，以及它们如何解决文档中的问题或挑战。

        请确保在每个步骤中，逐步推导出你的结论，并确保所有分析和建议具有逻辑连贯性。
        """
    }

def analyze_content(contents, api_key, prompt_template):
    if not api_key:
        st.error("请先输入有效的Google Gemini API密钥")
        return None
    
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-1.5-pro')
    
    parts = [prompt_template]
    for content in contents:
        if isinstance(content, dict) and content['type'] == 'image':
            image = Image.open(BytesIO(base64.b64decode(content['data'])))
            parts.extend([image, content['caption']])
        else:
            parts.append(str(content))
    
    try:
        response = model.generate_content(parts, generation_config=genai.types.GenerationConfig(
            temperature=temperature
        ))
        return response.text
    except Exception as e:
        st.error(f"生成内容时发生错误: {str(e)}")
        return None

prompt_options = get_prompt_options()
selected_option = st.selectbox("选择分析方式", list(prompt_options.keys()) + ["自定义提示"])
custom_prompt = st.text_area("输入自定义提示", height=200) if selected_option == "自定义提示" else prompt_options[selected_option]

uploaded_files = st.file_uploader("上传文件", type=["pdf", "pptx", "docx", "xlsx", "txt", "md", "jpg", "jpeg", "png", "gif", "bmp"], accept_multiple_files=True)

if uploaded_files:
    st.write(f"已上传 {len(uploaded_files)} 个文件")
    show_combined_content = st.checkbox("显示合并后的内容")
    
    if st.button("分析所有文件"):
        if not api_key:
            st.error("请先输入有效的Google Gemini API密钥")
        else:
            contents = []
            for file in uploaded_files:
                content = extract_content_from_file(file)
                contents.extend(content if isinstance(content, list) else [content])
            
            if show_combined_content:
                st.write("合并后的内容:")
                for content in contents:
                    if isinstance(content, dict) and content['type'] == 'image':
                        st.image(Image.open(BytesIO(base64.b64decode(content['data']))), caption=content['caption'])
                    else:
                        st.write(content)
            
            analysis_result = analyze_content(contents, api_key, custom_prompt)
            if analysis_result:
                st.session_state.analysis_result = analysis_result
                st.session_state.combined_content = "内容已处理,包括文本和图像"
                
                st.write("分析结果:")
                st.write(analysis_result)
                
                st.download_button(
                    label="下载分析结果",
                    data=analysis_result.encode('utf-8'),
                    file_name="analysis_result.txt",
                    mime="text/plain"
                )

st.write("使用说明:")
st.write("1. 输入您的Google Gemini API密钥")
st.write("2. 上传文件 (支持PDF、PPTX、DOCX、XLSX、TXT、MD、JPG、JPEG、PNG、GIF和BMP格式)")
st.write("3. 选择分析方式")
st.write("4. 点击'分析所有文件'按钮")
st.write("5. 查看分析结果")

st.markdown("---")
st.write("使用Gemini-1.5-Flash-Latest API进行文档内容分析")
